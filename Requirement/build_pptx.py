"""
HSNT WMS Board Presentation Builder
Theme: White + Light Green, Interactive Navigation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ─── Color Palette ────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN= RGBColor(0xE8, 0xF5, 0xEE)   # #E8F5EE
MID_GREEN  = RGBColor(0x4C, 0xAF, 0x7D)   # #4CAF7D
DARK_GREEN = RGBColor(0x1B, 0x6B, 0x3A)   # #1B6B3A
DEEP_GREEN = RGBColor(0x0D, 0x4A, 0x25)   # #0D4A25
ACCENT_GREEN=RGBColor(0x00, 0xC2, 0x7A)   # #00C27A
TEXT_DARK  = RGBColor(0x1A, 0x2E, 0x1A)   # #1A2E1A
TEXT_MID   = RGBColor(0x4A, 0x7A, 0x5A)   # #4A7A5A
TEXT_LIGHT = RGBColor(0x8A, 0xB8, 0x9A)   # #8AB89A
CARD_BG    = RGBColor(0xF4, 0xFB, 0xF6)   # #F4FBF6
BORDER_G   = RGBColor(0xB7, 0xE4, 0xC7)   # #B7E4C7
GRAY_MID   = RGBColor(0x64, 0x74, 0x80)   # #647480
OFF_WHITE  = RGBColor(0xF8, 0xFF, 0xFA)   # #F8FFFA

# ─── Slide size 16x9 (10" x 5.625") ──────────────────────────────────────────
W = Inches(10)
H = Inches(5.625)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # Completely blank

# ─── Helpers ──────────────────────────────────────────────────────────────────

def rgb_hex(r, g, b):
    return f"{r:02X}{g:02X}{b:02X}"

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None, radius=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=14, color=TEXT_DARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, valign=None, font_name="Calibri",
             bg_color=None, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    else:
        txBox.fill.background()
    txBox.line.fill.background()
    return txBox

def add_hyperlink_text(slide, text, x, y, w, h, target_slide_idx,
                       font_size=12, color=DARK_GREEN, bold=False, font_name="Calibri",
                       align=PP_ALIGN.LEFT):
    """Add text with internal slide hyperlink using OOXML rPr"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = False
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = align

    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.underline = True

    # Add internal slide hyperlink via relationship
    slide_part = slide.part
    target_slide = prs.slides[target_slide_idx]
    rId = slide_part.relate_to(target_slide.part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")

    r_elem = run._r
    rPr = r_elem.get_or_add_rPr()

    hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)

    txBox.fill.background()
    txBox.line.fill.background()
    return txBox

def add_nav_button(slide, label, x, y, target_slide_idx):
    """Small green button linking to a slide"""
    btn = add_rect(slide, x, y, 1.2, 0.32, LIGHT_GREEN, MID_GREEN, 0.75)
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.2), Inches(0.32))
    txBox.word_wrap = False
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.color.rgb = DARK_GREEN
    run.font.bold = True

    slide_part = slide.part
    target_slide = prs.slides[target_slide_idx]
    rId = slide_part.relate_to(target_slide.part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
    r_elem = run._r
    rPr = r_elem.get_or_add_rPr()
    hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    txBox.fill.background()
    txBox.line.fill.background()
    return btn

def add_slide_header(slide, title, subtitle=None):
    """Standard green header bar for content slides"""
    add_rect(slide, 0, 0, 10, 0.85, DARK_GREEN)
    add_rect(slide, 0, 0.85, 10, 0.055, MID_GREEN)
    add_text(slide, title, 0.35, 0.08, 8.5, 0.65,
             font_size=24, color=WHITE, bold=True, valign="middle", font_name="Calibri")
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.08, 9.5, 0.65,
                 font_size=10, color=BORDER_G, italic=True,
                 align=PP_ALIGN.RIGHT, font_name="Calibri")

def add_card(slide, x, y, w, h, title=None, body_lines=None,
             title_color=DARK_GREEN, body_color=TEXT_DARK,
             card_bg=CARD_BG, border_color=BORDER_G):
    """Rounded-look card with optional title and bullets"""
    add_rect(slide, x, y, w, h, card_bg, border_color, 0.6)
    ty = y + 0.1
    if title:
        add_text(slide, title, x + 0.15, ty, w - 0.3, 0.32,
                 font_size=11, color=title_color, bold=True, font_name="Calibri")
        ty += 0.33
    if body_lines:
        for line in body_lines:
            add_text(slide, f"  {line}", x + 0.1, ty, w - 0.2, 0.25,
                     font_size=9.5, color=body_color, font_name="Calibri")
            ty += 0.25

def slide_footer(slide, slide_number, total=11):
    """Slide number footer"""
    add_text(slide, f"{slide_number} / {total}", 9.3, 5.35, 0.65, 0.22,
             font_size=8, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT, font_name="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
s1.background.fill.solid()
s1.background.fill.fore_color.rgb = WHITE

# Deep green left panel
add_rect(s1, 0, 0, 4.2, 5.625, DEEP_GREEN)
add_rect(s1, 4.2, 0, 0.08, 5.625, MID_GREEN)

# Company tag
add_text(s1, "HSNT", 0.3, 0.35, 3.5, 0.7,
         font_size=32, color=ACCENT_GREEN, bold=True, font_name="Calibri")
add_text(s1, "Highpoint Service Network Thailand", 0.3, 0.95, 3.5, 0.4,
         font_size=9, color=BORDER_G, font_name="Calibri")

# Divider line
add_rect(s1, 0.3, 1.42, 3.0, 0.04, BORDER_G)

add_text(s1, "WMS", 0.3, 1.55, 3.5, 1.1,
         font_size=72, color=WHITE, bold=True, font_name="Calibri")
add_text(s1, "Warehouse Management System", 0.3, 2.6, 3.5, 0.5,
         font_size=13, color=LIGHT_GREEN, font_name="Calibri")

add_text(s1, "นำเสนอต่อคณะกรรมการบริหาร", 0.3, 3.3, 3.5, 0.4,
         font_size=10, color=TEXT_LIGHT, italic=True, font_name="Calibri")
add_text(s1, "2026", 0.3, 3.65, 3.5, 0.5,
         font_size=14, color=BORDER_G, font_name="Calibri")

# Right side content summary
add_text(s1, "ระบบบริหารจัดการคลังสินค้า", 4.5, 0.6, 5.2, 0.6,
         font_size=20, color=DARK_GREEN, bold=True, font_name="Calibri")
add_text(s1, "ครบวงจร | ทันสมัย | ปลอดภัย", 4.5, 1.15, 5.2, 0.4,
         font_size=12, color=MID_GREEN, font_name="Calibri")

stats = [
    ("4", "Portals หลัก"),
    ("12", "User Roles"),
    ("5+", "Key Workflows"),
    ("20+", "KPI Dashboard"),
]
sx = [4.5, 6.2, 7.9, 9.0]
for i, (num, lbl) in enumerate(stats):
    cx = 4.5 + i * 1.45
    if cx + 1.3 > 9.9:
        break
    add_rect(s1, cx, 1.75, 1.35, 1.1, LIGHT_GREEN, BORDER_G, 0.5)
    add_text(s1, num, cx, 1.8, 1.35, 0.55,
             font_size=30, color=DARK_GREEN, bold=True,
             align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(s1, lbl, cx, 2.35, 1.35, 0.4,
             font_size=8.5, color=TEXT_MID, align=PP_ALIGN.CENTER, font_name="Calibri")

add_text(s1, "คลิกเพื่อเริ่มต้น ›", 4.5, 5.1, 5.2, 0.42,
         font_size=10, color=TEXT_LIGHT, italic=True, font_name="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda / Interactive TOC
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
s2.background.fill.solid()
s2.background.fill.fore_color.rgb = WHITE
add_slide_header(s2, "สารบัญ — Agenda", "คลิกหัวข้อเพื่อไปยังสไลด์นั้น")

add_text(s2, "Interactive Navigation — คลิกหัวข้อด้านล่างเพื่อไปยังเนื้อหา",
         0.4, 0.98, 9.2, 0.35, font_size=10, color=TEXT_MID, italic=True, font_name="Calibri")

# Will add hyperlinks after all slides are created
# Placeholder markers stored for later
agenda_items = [
    (2,  "01", "ที่มาและวัตถุประสงค์",        "Overview & Objectives"),
    (3,  "02", "โครงสร้างระบบ 4 Portals",    "System Architecture"),
    (4,  "03", "12 บทบาทผู้ใช้งาน",           "User Roles & Permissions"),
    (5,  "04", "กระบวนการทำงานหลัก",          "Key Workflows"),
    (6,  "05", "การควบคุมสถานะสต็อก",        "Inventory Status Control"),
    (7,  "06", "Dashboard ผู้บริหาร",         "Executive KPI Dashboard"),
    (8,  "07", "สถานที่ปฏิบัติการ",            "Operational Sites"),
    (9,  "08", "เทคโนโลยีและการออกแบบ",      "Technology & Design"),
    (10, "09", "แผนพัฒนาในอนาคต",            "Future Roadmap"),
    (10, "10", "สรุปสำหรับคณะกรรมการ",       "Board Summary"),
]

# Layout 2 columns x 5 rows
col_positions = [0.35, 5.1]
row_start = 1.4
row_h = 0.72
for idx, (slide_idx, num, th, en) in enumerate(agenda_items):
    col = idx % 2
    row = idx // 2
    cx = col_positions[col]
    cy = row_start + row * row_h

    add_rect(s2, cx, cy, 4.5, 0.6, LIGHT_GREEN, BORDER_G, 0.5)
    # Number badge
    add_rect(s2, cx, cy, 0.38, 0.6, MID_GREEN)
    add_text(s2, num, cx, cy + 0.1, 0.38, 0.4,
             font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")

slide_footer(s2, 2)

# We will link TOC items after all slides are created
toc_items_data = agenda_items  # save for later

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Overview & Objectives
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
s3.background.fill.solid()
s3.background.fill.fore_color.rgb = WHITE
add_slide_header(s3, "01  ที่มาและวัตถุประสงค์", "Overview & Objectives")

add_text(s3, "ปัญหาที่ต้องการแก้ไข", 0.4, 1.0, 4.5, 0.35,
         font_size=13, color=DARK_GREEN, bold=True, font_name="Calibri")
problems = [
    "การบริหารสต็อกยังใช้ Excel / Manual",
    "ข้อมูลสต็อกไม่ Real-time, เกิด Error บ่อย",
    "ไม่มีระบบ Approval Workflow ที่ชัดเจน",
    "ขาดการ Track DOA / RTV อย่างเป็นระบบ",
    "ผู้บริหารขาด Dashboard ติดตาม KPI",
]
for i, p in enumerate(problems):
    add_rect(s3, 0.4, 1.38 + i*0.56, 0.32, 0.32, MID_GREEN)
    add_text(s3, "!", 0.4, 1.38 + i*0.56, 0.32, 0.32,
             font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(s3, p, 0.8, 1.4 + i*0.56, 3.9, 0.3,
             font_size=10, color=TEXT_DARK, font_name="Calibri")

# Arrow
add_rect(s3, 4.7, 2.3, 0.5, 0.3, LIGHT_GREEN)
add_text(s3, "›", 4.7, 2.28, 0.5, 0.3,
         font_size=18, color=MID_GREEN, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")

add_text(s3, "วัตถุประสงค์ WMS", 5.4, 1.0, 4.2, 0.35,
         font_size=13, color=DARK_GREEN, bold=True, font_name="Calibri")
goals = [
    ("🏭", "Real-time Stock Visibility ทุก Location"),
    ("✅", "Approval Workflow อัตโนมัติ"),
    ("📦", "Pick / Pack / Ship ครบวงจร"),
    ("🔄", "RMA & RTV Management"),
    ("📊", "Executive Dashboard & KPI"),
]
for i, (icon, goal) in enumerate(goals):
    add_rect(s3, 5.4, 1.38 + i*0.56, 4.2, 0.42, CARD_BG, BORDER_G, 0.5)
    add_text(s3, icon, 5.45, 1.4 + i*0.56, 0.4, 0.3,
             font_size=12, color=MID_GREEN, font_name="Calibri")
    add_text(s3, goal, 5.85, 1.42 + i*0.56, 3.65, 0.3,
             font_size=10, color=TEXT_DARK, font_name="Calibri")

add_rect(s3, 0, 5.25, 10, 0.375, LIGHT_GREEN)
add_text(s3, "HSNT WMS คือระบบที่เปลี่ยน Manual Process ไปสู่ Digital Workflow ที่ตรวจสอบได้, มีประสิทธิภาพ และรองรับการเติบโต",
         0.4, 5.27, 9.2, 0.33,
         font_size=9.5, color=DARK_GREEN, italic=True, font_name="Calibri")
slide_footer(s3, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture: 4 Portals
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
s4.background.fill.solid()
s4.background.fill.fore_color.rgb = WHITE
add_slide_header(s4, "02  โครงสร้างระบบ — 4 Portals", "System Architecture")

portals = [
    ("Admin\nSystem",       "⚙",  DEEP_GREEN,  WHITE,
     ["User & Role Mgmt", "Product / Brand / Vendor Master",
      "Warehouse & Rack Master", "Approval Workflow Setup", "Audit Log & Notification"]),
    ("Warehouse\nManagement","🏭", DARK_GREEN,  WHITE,
     ["Goods Receiving", "Putaway & Stock Transfer",
      "Pick / Pack / Ship", "Cycle Count", "Warehouse Dashboard"]),
    ("Requester\nPortal",   "📋", MID_GREEN,   WHITE,
     ["Create Withdrawal Request", "Track Request Status",
      "Confirm RMA Usage", "Declare DOA / Defective", "Unused Goods Return"]),
    ("RTV\nManagement",     "🔄", ACCENT_GREEN, WHITE,
     ["DOA & Defective Review", "RTV Case Creation & Approval",
      "Vendor Return Shipment", "Replacement Tracking", "Credit Note & Closure"]),
]

px = [0.25, 2.75, 5.25, 7.75]
for i, (name, icon, bg, fg, items) in enumerate(portals):
    x = px[i]
    add_rect(s4, x, 1.0, 2.3, 4.35, CARD_BG, BORDER_G, 0.6)
    add_rect(s4, x, 1.0, 2.3, 1.05, bg)
    add_text(s4, icon, x, 1.02, 2.3, 0.45,
             font_size=20, color=fg, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(s4, name, x, 1.44, 2.3, 0.55,
             font_size=10.5, color=fg, bold=True,
             align=PP_ALIGN.CENTER, font_name="Calibri")
    for j, item in enumerate(items):
        add_text(s4, f"• {item}", x + 0.12, 2.12 + j * 0.46, 2.1, 0.38,
                 font_size=8.5, color=TEXT_DARK, font_name="Calibri")

slide_footer(s4, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — User Roles
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
s5.background.fill.solid()
s5.background.fill.fore_color.rgb = WHITE
add_slide_header(s5, "03  บทบาทผู้ใช้งาน 12 Roles", "User Roles & Access Control")

roles = [
    ("01", "System Administrator",   "Full system access & configuration"),
    ("02", "Warehouse Manager",      "Oversee all warehouse operations"),
    ("03", "Warehouse Supervisor",   "Supervise daily warehouse tasks"),
    ("04", "Warehouse Staff",        "Execute pick / pack / receive / ship"),
    ("05", "Requester",              "Create & track withdrawal requests"),
    ("06", "Department Approver",    "Approve requests for own department"),
    ("07", "RMA / Service Team",     "Confirm RMA usage & item status"),
    ("08", "RTV Officer",            "Process DOA & vendor returns"),
    ("09", "Finance / Acct Viewer",  "View stock values & transactions"),
    ("10", "Brand / Vendor Viewer",  "View brand-specific stock & RTV"),
    ("11", "Management Viewer",      "Executive dashboard & reports"),
    ("12", "Auditor",                "Read-only audit log & cycle count"),
]

colors_roles = [DEEP_GREEN, DARK_GREEN, DARK_GREEN, MID_GREEN,
                MID_GREEN, ACCENT_GREEN, ACCENT_GREEN, RGBColor(0x07,0x9A,0x60),
                GRAY_MID, GRAY_MID, RGBColor(0x1B,0x6B,0x3A), RGBColor(0x4A,0x7A,0x5A)]

cols = 3
col_w = 3.1
row_h = 0.56
start_x = [0.25, 3.45, 6.65]
start_y = 1.05
for i, (num, role, desc) in enumerate(roles):
    col = i % cols
    row = i // cols
    cx = start_x[col]
    cy = start_y + row * row_h
    c = colors_roles[i]
    add_rect(s5, cx, cy, col_w, 0.5, LIGHT_GREEN, BORDER_G, 0.5)
    add_rect(s5, cx, cy, 0.42, 0.5, c)
    add_text(s5, num, cx, cy + 0.08, 0.42, 0.32,
             font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(s5, role, cx + 0.48, cy + 0.04, col_w - 0.55, 0.25,
             font_size=9.5, color=DARK_GREEN, bold=True, font_name="Calibri")
    add_text(s5, desc, cx + 0.48, cy + 0.26, col_w - 0.55, 0.22,
             font_size=8, color=TEXT_MID, font_name="Calibri")

slide_footer(s5, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Key Workflows
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
s6.background.fill.solid()
s6.background.fill.fore_color.rgb = WHITE
add_slide_header(s6, "04  กระบวนการทำงานหลัก", "Key Workflows")

flows = [
    ("Goods Receiving",      ["สินค้าเข้า","ตรวจสอบ","สร้าง Receiving","Putaway","สต็อกพร้อมใช้"]),
    ("Withdrawal Request",   ["Requester ขอเบิก","Validation","Approval","Reserve Stock","Picking Task"]),
    ("Pick/Pack/Ship",       ["รับ Task","Pick ตาม Location","Pack & Label","Ship","Receiver ยืนยัน"]),
    ("RMA Usage Confirm",    ["เบิกสินค้าออก","ใช้งาน","แจ้งผลการใช้","DOA/คืน/ใช้แล้ว","ปิด Case"]),
    ("RTV Process",          ["แจ้ง DOA","RTV Review","Approve","คืน Vendor","ติดตามผล & ปิด"]),
]

for fi, (fname, steps) in enumerate(flows):
    fy = 1.0 + fi * 0.89
    add_rect(s6, 0.2, fy, 1.5, 0.72, LIGHT_GREEN, BORDER_G, 0.5)
    add_text(s6, fname, 0.25, fy + 0.08, 1.4, 0.55,
             font_size=8.5, color=DARK_GREEN, bold=True, font_name="Calibri")

    for si, step in enumerate(steps):
        sx = 1.85 + si * 1.66
        sc = MID_GREEN if si == 0 else (DARK_GREEN if si == len(steps)-1 else CARD_BG)
        st = WHITE if si in (0, len(steps)-1) else TEXT_DARK
        add_rect(s6, sx, fy + 0.08, 1.5, 0.5, sc, BORDER_G, 0.5)
        add_text(s6, step, sx + 0.05, fy + 0.1, 1.4, 0.45,
                 font_size=8, color=st, align=PP_ALIGN.CENTER, font_name="Calibri")
        if si < len(steps) - 1:
            add_text(s6, "›", sx + 1.52, fy + 0.18, 0.2, 0.25,
                     font_size=11, color=MID_GREEN, bold=True, font_name="Calibri")

slide_footer(s6, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Inventory Status
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
s7.background.fill.solid()
s7.background.fill.fore_color.rgb = WHITE
add_slide_header(s7, "05  การควบคุมสถานะสต็อก", "Inventory Status Management")

statuses = [
    ("Pending Receiving",  "EFF7F2", "4CAF7D"),
    ("Pending Inspection", "EFF7F2", "4CAF7D"),
    ("Available",          "D4EDDA", "1B6B3A"),
    ("Reserved",           "D4EDDA", "1B6B3A"),
    ("Picking",            "FFF3CD", "856404"),
    ("Packed",             "FFF3CD", "856404"),
    ("Shipped",            "CCE5FF", "004085"),
    ("Issued to RMA",      "CCE5FF", "004085"),
    ("Consumed",           "D6D8DB", "383D41"),
    ("Quarantine",         "F8D7DA", "721C24"),
    ("DOA",                "F8D7DA", "721C24"),
    ("RTV Pending",        "F8D7DA", "721C24"),
    ("RTV Shipped",        "E2CFEE", "4A235A"),
    ("Returned Unused",    "EFF7F2", "4CAF7D"),
    ("Closed",             "D6D8DB", "383D41"),
    ("Cancelled",          "D6D8DB", "383D41"),
]

ncols = 4
cw = 2.3
rh = 0.46
sx0 = 0.35
sy0 = 1.0
for i, (label, bg_hex, tx_hex) in enumerate(statuses):
    col = i % ncols
    row = i // ncols
    cx = sx0 + col * (cw + 0.13)
    cy = sy0 + row * rh

    r, g, b = int(bg_hex[0:2],16), int(bg_hex[2:4],16), int(bg_hex[4:6],16)
    tr, tg, tb = int(tx_hex[0:2],16), int(tx_hex[2:4],16), int(tx_hex[4:6],16)
    add_rect(s7, cx, cy, cw, 0.38, RGBColor(r,g,b), RGBColor(tr,tg,tb), 0.5)
    add_text(s7, f"● {label}", cx + 0.1, cy + 0.04, cw - 0.15, 0.3,
             font_size=9, color=RGBColor(tr,tg,tb), bold=True, font_name="Calibri")

add_text(s7, "● สีเขียว = พร้อมใช้  ● สีเหลือง = ระหว่างดำเนินการ  ● สีแดง = ปัญหา / กักกัน  ● สีเทา = ปิด",
         0.35, 5.25, 9.3, 0.33,
         font_size=9, color=TEXT_MID, align=PP_ALIGN.CENTER, font_name="Calibri")
slide_footer(s7, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Executive Dashboard
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
s8.background.fill.solid()
s8.background.fill.fore_color.rgb = WHITE
add_slide_header(s8, "06  Executive Dashboard & KPI", "Real-time Monitoring for Management")

kpis = [
    ("มูลค่าสต็อกรวม",   "฿ XXM",  DARK_GREEN),
    ("สินค้า Available",  "XX,XXX", MID_GREEN),
    ("DOA รอดำเนินการ",   "XX",     RGBColor(0xC0,0x39,0x2B)),
    ("RTV รอส่งคืน",      "XX",     RGBColor(0xE6,0x7E,0x22)),
    ("Warehouse Util.",   "XX%",    MID_GREEN),
    ("SLA Completion",    "XX%",    DARK_GREEN),
]

kx = [0.3, 2.05, 3.8, 5.55, 7.3, 8.7]
kw = 1.6
for i, (label, val, c) in enumerate(kpis):
    cx = 0.3 + i * 1.66
    add_rect(s8, cx, 1.0, kw, 1.05, LIGHT_GREEN, BORDER_G, 0.6)
    add_text(s8, val, cx, 1.05, kw, 0.55,
             font_size=24, color=c, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(s8, label, cx, 1.58, kw, 0.38,
             font_size=8, color=TEXT_MID, align=PP_ALIGN.CENTER, font_name="Calibri")

# Charts placeholders
charts_info = [
    (0.3, 2.22, 4.5, 2.85, "สต็อกแยกตาม Brand & Warehouse",
     [("Brand A", 35), ("Brand B", 28), ("Brand C", 22), ("Other", 15)]),
    (5.0, 2.22, 4.7, 2.85, "Monthly Withdrawal Volume (6 เดือน)",
     None),
]

# Left: donut-like pie chart simulation with bars
add_rect(s8, 0.3, 2.22, 4.5, 2.85, CARD_BG, BORDER_G, 0.5)
add_text(s8, "สต็อกแยกตาม Brand", 0.5, 2.3, 4.1, 0.35,
         font_size=11, color=DARK_GREEN, bold=True, font_name="Calibri")

brands = [("Brand A", 35, MID_GREEN), ("Brand B", 28, DARK_GREEN),
          ("Brand C", 22, ACCENT_GREEN), ("Other", 15, RGBColor(0xB7,0xE4,0xC7))]
bar_y = 2.72
for bi, (bn, pct, bc) in enumerate(brands):
    bw = (pct / 100) * 3.5
    add_rect(s8, 0.5, bar_y + bi*0.48, bw, 0.32, bc)
    add_text(s8, f"{bn}  {pct}%", 0.5, bar_y + bi*0.48, 3.5, 0.3,
             font_size=9, color=WHITE, bold=True, font_name="Calibri")

# Right: inventory aging
add_rect(s8, 5.0, 2.22, 4.7, 2.85, CARD_BG, BORDER_G, 0.5)
add_text(s8, "Inventory Aging & SLA", 5.2, 2.3, 4.3, 0.35,
         font_size=11, color=DARK_GREEN, bold=True, font_name="Calibri")

aging = [("< 30 วัน",  "สดใหม่",   72, MID_GREEN),
         ("30-90 วัน", "ปกติ",     18, RGBColor(0xF0,0xB4,0x29)),
         ("> 90 วัน",  "ตรวจสอบ",  10, RGBColor(0xC0,0x39,0x2B))]
for ai, (rng, lbl, pct, ac) in enumerate(aging):
    ay = 2.72 + ai * 0.56
    bw = (pct / 100) * 3.6
    add_rect(s8, 5.2, ay, bw, 0.35, ac)
    add_text(s8, f"{rng}  ({lbl})  {pct}%", 5.2, ay, 3.6, 0.32,
             font_size=8.5, color=WHITE, bold=True, font_name="Calibri")

slide_footer(s8, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Operational Sites
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
s9.background.fill.solid()
s9.background.fill.fore_color.rgb = WHITE
add_slide_header(s9, "07  สถานที่ปฏิบัติการ", "Operational Warehouse Sites")

sites = [
    ("Tower B1FL",     "คลังชั้น B1 อาคาร Tower",    "Spare Parts, Consignment"),
    ("Tower B2FL",     "คลังชั้น B2 อาคาร Tower",    "Finished Goods, RMA Buffer"),
    ("Tower C1FL",     "คลังชั้น C1 อาคาร Tower",    "Consignment, Overflow"),
    ("SCG Warehouse",  "คลังสินค้า SCG (Off-site)",   "Bulk / Long-term Storage"),
    ("RTV / Quarantine","พื้นที่กักกันสินค้า",        "DOA, Defective, RTV Pending"),
]

for i, (name, desc, type_) in enumerate(sites):
    cx = 0.3 + (i % 3) * 3.25
    cy = 1.1 + (i // 3) * 2.05
    add_rect(s9, cx, cy, 3.0, 1.75, LIGHT_GREEN, BORDER_G, 0.6)
    add_rect(s9, cx, cy, 3.0, 0.5, DARK_GREEN)
    add_text(s9, f"📍 {name}", cx + 0.1, cy + 0.06, 2.8, 0.38,
             font_size=11, color=WHITE, bold=True, font_name="Calibri")
    add_text(s9, desc, cx + 0.1, cy + 0.56, 2.8, 0.38,
             font_size=9.5, color=TEXT_DARK, font_name="Calibri")
    add_text(s9, f"Type: {type_}", cx + 0.1, cy + 0.92, 2.8, 0.32,
             font_size=8.5, color=TEXT_MID, italic=True, font_name="Calibri")
    add_text(s9, "● Real-time Stock Tracking\n● Location / Rack / Bin Mgmt",
             cx + 0.1, cy + 1.22, 2.8, 0.45,
             font_size=8, color=MID_GREEN, font_name="Calibri")

slide_footer(s9, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Technology & Design
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
s10.background.fill.solid()
s10.background.fill.fore_color.rgb = WHITE
add_slide_header(s10, "08  เทคโนโลยีและการออกแบบ", "Technology & System Design")

add_text(s10, "UI/UX Design Principles", 0.4, 1.0, 4.5, 0.38,
         font_size=13, color=DARK_GREEN, bold=True, font_name="Calibri")
ui_items = [
    "Modern SaaS Dashboard Style",
    "Role-Based Menu & Access Control",
    "Responsive Layout (Desktop / Tablet)",
    "Notification Center แบบ Real-time",
    "Audit Trail ทุก Transaction",
    "Approval Workflow Engine",
    "Status Badge & Color Coding",
    "Kanban Board: Pick/Pack/Ship",
]
for i, item in enumerate(ui_items):
    add_rect(s10, 0.4, 1.45 + i*0.47, 0.28, 0.28, MID_GREEN)
    add_text(s10, "✓", 0.4, 1.46 + i*0.47, 0.28, 0.26,
             font_size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(s10, item, 0.76, 1.47 + i*0.47, 3.7, 0.27,
             font_size=9.5, color=TEXT_DARK, font_name="Calibri")

# Right column: tech stack / integration
add_text(s10, "การออกแบบเชิงระบบ", 5.4, 1.0, 4.2, 0.38,
         font_size=13, color=DARK_GREEN, bold=True, font_name="Calibri")
tech = [
    ("🔒", "Role-Based Access Control (RBAC)"),
    ("📝", "Audit Log ทุก Action"),
    ("🔔", "Notification Engine"),
    ("📊", "Executive Dashboard Real-time"),
    ("📱", "รองรับ Mobile App (Future)"),
    ("🔗", "API Integration Ready"),
    ("📈", "Power BI / Analytics Ready"),
    ("🏷",  "Barcode / QR Scan Ready"),
]
for i, (icon, desc) in enumerate(tech):
    add_rect(s10, 5.4, 1.45 + i*0.47, 4.2, 0.38, CARD_BG, BORDER_G, 0.5)
    add_text(s10, icon, 5.45, 1.47 + i*0.47, 0.38, 0.28,
             font_size=12, font_name="Calibri")
    add_text(s10, desc, 5.85, 1.49 + i*0.47, 3.65, 0.28,
             font_size=9.5, color=TEXT_DARK, font_name="Calibri")

slide_footer(s10, 10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Future Roadmap
# ══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
s11.background.fill.solid()
s11.background.fill.fore_color.rgb = WHITE
add_slide_header(s11, "09  แผนพัฒนาในอนาคต", "Future Roadmap")

phases = [
    ("Phase 1\nNow → Q3",    "Core WMS Launch",       DEEP_GREEN,
     ["Admin, Warehouse, Requester, RTV Portals",
      "12 User Roles & RBAC",
      "Full Workflow: Receive → Ship",
      "Executive Dashboard"]),
    ("Phase 2\nQ4 → Q1+1",  "Enhanced Integration",   DARK_GREEN,
     ["Barcode / QR Code Scanner",
      "Mobile Warehouse App",
      "RMA System Integration",
      "Vendor Portal (Read-only)"]),
    ("Phase 3\nQ2+1 →",     "Analytics & AI",         MID_GREEN,
     ["Power BI Dashboard Integration",
      "Inventory Forecasting",
      "DOA / RTV Trend Analytics",
      "Cycle Count Automation",
      "SLA Monitoring & Alert"]),
]

pw = 2.9
for i, (phase, title, c, items) in enumerate(phases):
    px2 = 0.35 + i * 3.2
    add_rect(s11, px2, 1.0, pw, 4.35, CARD_BG, BORDER_G, 0.5)
    add_rect(s11, px2, 1.0, pw, 0.72, c)
    add_text(s11, phase, px2, 1.02, pw, 0.34,
             font_size=9, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(s11, title, px2, 1.35, pw, 0.32,
             font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")

    for ji, item in enumerate(items):
        add_text(s11, f"▸  {item}", px2 + 0.1, 1.82 + ji * 0.56, pw - 0.2, 0.48,
                 font_size=9, color=TEXT_DARK, font_name="Calibri")

# Timeline bar
add_rect(s11, 0.35, 5.25, 9.3, 0.2, LIGHT_GREEN, BORDER_G, 0.5)
for ti, (lbl, cx) in enumerate([("Now", 0.35), ("Q3", 3.55), ("Q4", 6.75), ("Q2+1", 9.4)]):
    add_text(s11, lbl, cx - 0.2, 5.26, 0.8, 0.18,
             font_size=8, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")

slide_footer(s11, 11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Board Summary
# ══════════════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(blank_layout)
s12.background.fill.solid()
s12.background.fill.fore_color.rgb = WHITE

add_rect(s12, 0, 0, 10, 1.1, DEEP_GREEN)
add_rect(s12, 0, 1.1, 10, 0.055, MID_GREEN)
add_text(s12, "10  สรุปสำหรับคณะกรรมการบริหาร", 0.4, 0.12, 9, 0.65,
         font_size=24, color=WHITE, bold=True, font_name="Calibri")

add_text(s12, "HSNT WMS — ระบบบริหารจัดการคลังสินค้าที่ครบวงจรและพร้อม Scale",
         0.5, 1.22, 9, 0.42,
         font_size=13, color=DARK_GREEN, bold=True, font_name="Calibri")

summary_points = [
    ("✅", "ครอบคลุม", "4 Portals ดูแลทุกมิติของ Warehouse Operations ตั้งแต่ Admin ถึง RTV"),
    ("✅", "ปลอดภัย",  "12 Roles พร้อม RBAC ทุก Action มี Audit Trail ตรวจสอบได้"),
    ("✅", "Real-time","Executive Dashboard ติดตาม KPI และ Stock ได้ทันที"),
    ("✅", "รองรับการเติบโต","Roadmap 3 Phase พร้อม Mobile App, Analytics และ AI ในอนาคต"),
]

for i, (icon, head, body) in enumerate(summary_points):
    cx = 0.4 + (i % 2) * 4.8
    cy = 1.75 + (i // 2) * 1.1
    add_rect(s12, cx, cy, 4.5, 0.95, LIGHT_GREEN, BORDER_G, 0.5)
    add_text(s12, icon, cx + 0.1, cy + 0.08, 0.4, 0.4,
             font_size=16, font_name="Calibri")
    add_text(s12, head, cx + 0.55, cy + 0.06, 3.8, 0.32,
             font_size=11.5, color=DARK_GREEN, bold=True, font_name="Calibri")
    add_text(s12, body, cx + 0.55, cy + 0.38, 3.8, 0.5,
             font_size=9, color=TEXT_DARK, font_name="Calibri")

add_rect(s12, 0, 4.85, 10, 0.775, DARK_GREEN)
add_text(s12, "HSNT WMS พร้อมยกระดับการบริหารคลังสินค้าจาก Manual สู่ Digital อย่างมีประสิทธิภาพ",
         0.5, 4.92, 9, 0.4,
         font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(s12, "Highpoint Service Network Thailand Co., Ltd.  |  2026",
         0.5, 5.26, 9, 0.3,
         font_size=9, color=BORDER_G, align=PP_ALIGN.CENTER, font_name="Calibri")

slide_footer(s12, 12, 12)

# ══════════════════════════════════════════════════════════════════════════════
# NOW ADD HYPERLINKS — all slides exist, we can link
# ══════════════════════════════════════════════════════════════════════════════

# TOC: Agenda items with clickable links
# Slide indices: 0-based in prs.slides
# s1=0, s2=1, s3=2, s4=3, s5=4, s6=5, s7=6, s8=7, s9=8, s10=9, s11=10, s12=11

toc_links = [
    (2, "01  ที่มาและวัตถุประสงค์",      "Overview & Objectives"),
    (3, "02  โครงสร้างระบบ 4 Portals", "System Architecture"),
    (4, "03  บทบาทผู้ใช้งาน 12 Roles", "User Roles & Permissions"),
    (5, "04  กระบวนการทำงานหลัก",       "Key Workflows"),
    (6, "05  การควบคุมสถานะสต็อก",     "Inventory Status Control"),
    (7, "06  Dashboard ผู้บริหาร",      "Executive KPI Dashboard"),
    (8, "07  สถานที่ปฏิบัติการ",         "Operational Sites"),
    (9, "08  เทคโนโลยีและการออกแบบ",   "Technology & Design"),
    (10,"09  แผนพัฒนาในอนาคต",         "Future Roadmap"),
    (11,"10  สรุปสำหรับคณะกรรมการ",    "Board Summary"),
]

col_positions = [0.35, 5.1]
row_start = 1.4
row_h = 0.72

for idx, (slide_idx, th, en) in enumerate(toc_links):
    col = idx % 2
    row = idx // 2
    cx = col_positions[col]
    cy = row_start + row * row_h
    # Add clickable text over the card
    add_hyperlink_text(s2, f"  {th}", cx + 0.42, cy + 0.04, 3.95, 0.3,
                       slide_idx, font_size=10.5, color=DARK_GREEN, bold=True)
    add_hyperlink_text(s2, f"  {en}", cx + 0.42, cy + 0.34, 3.95, 0.22,
                       slide_idx, font_size=8.5, color=TEXT_MID, bold=False)

# Add "◀ Back to Agenda" button on slides 3-12
content_slides = [s3, s4, s5, s6, s7, s8, s9, s10, s11, s12]
for cs in content_slides:
    add_nav_button(cs, "◀ เมนูหลัก", 8.6, 5.18, 1)  # link to slide index 1 = s2

# Add "Next ›" on title slide going to agenda
add_nav_button(s1, "เริ่มต้น ›", 8.6, 5.18, 1)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
output_path = r"D:\WMS_Project\Requirement\HSNT_WMS_Board_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
